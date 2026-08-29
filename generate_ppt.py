from pptx import Presentation
from pptx.util import Inches, Pt

# Create presentation
prs = Presentation()

# ==========================================
# Slide 1: PROBLEM UNDERSTANDING AND OBJECTIVE
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "1. PROBLEM UNDERSTANDING AND OBJECTIVE"
tf = slide.placeholders[1].text_frame

p = tf.add_paragraph()
p.text = "Objective: Build a Proof-of-Concept for Firstsource that intelligently extracts, structures, and audits unstructured clinical documents."
p.space_after = Pt(14)

p = tf.add_paragraph()
p.text = "The Problem Statement:"
p = tf.add_paragraph()
p.text = "- Healthcare staff spend countless hours manually extracting data from fragmented documents (lab reports, handwritten notes, intake forms)."
p.level = 1
p = tf.add_paragraph()
p.text = "- This manual process is slow, expensive, and highly prone to human error, leading to billing mistakes and delayed patient care."
p.level = 1


# ==========================================
# Slide 2: SOLUTION ARCHITECTURE AND DESIGN FLOW
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "2. SOLUTION ARCHITECTURE AND DESIGN FLOW"
tf = slide.placeholders[1].text_frame

p = tf.add_paragraph()
p.text = "Workflow & Data Flow:"
p = tf.add_paragraph()
p.text = "1. Ingestion: Users upload fragmented files (PDFs, JPGs, TXTs) via the frontend."
p.level = 1
p = tf.add_paragraph()
p.text = "2. Processing: FastAPI backend parses documents into image matrices."
p.level = 1
p = tf.add_paragraph()
p.text = "3. AI Extraction: Dual-LLM Ensemble (Gemini Flash & Groq Llama Vision) process data in parallel."
p.level = 1
p = tf.add_paragraph()
p.text = "4. Data Storage: Extracted JSON and raw PDF binaries are stored in MongoDB GridFS."
p.level = 1
p = tf.add_paragraph()
p.text = "5. Visualization: A split-screen dashboard renders the native document alongside the AI insights."
p.level = 1


# ==========================================
# Slide 3: IMPLEMENTATION HIGHLIGHTS
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "3. IMPLEMENTATION HIGHLIGHTS"
tf = slide.placeholders[1].text_frame

p = tf.add_paragraph()
p.text = "Technical Decisions & AI Logic:"
p = tf.add_paragraph()
p.text = "- Single-Pass Prompting: Collapsed multi-pass LLM extraction into a single, highly efficient extraction + scoring prompt to cut latency by 50%."
p.level = 1
p = tf.add_paragraph()
p.text = "- Dynamic Ensemble Merging: Both LLM outputs are mathematically merged. If models disagree, the system flags the field for human review (Safety First)."
p.level = 1
p = tf.add_paragraph()
p.text = "- UI/UX: Built a bespoke, glassmorphism dark-mode interface entirely in Vanilla JS/CSS for a premium, lightweight user experience without heavy frameworks."
p.level = 1


# ==========================================
# Slide 4: CHALLENGES AND LEARNINGS
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "4. CHALLENGES AND LEARNINGS"
tf = slide.placeholders[1].text_frame

p = tf.add_paragraph()
p.text = "Key Challenges & Trade-offs:"
p = tf.add_paragraph()
p.text = "- API Rate Limits: Free-tier LLMs (Gemini/Groq) strictly limit tokens and image counts. We implemented dynamic fallbacks and a hard 3-page cap to prevent server crashes."
p.level = 1
p = tf.add_paragraph()
p.text = "- JSON Truncation: Complex medical docs exceeded standard token limits, crashing the parser. Fixed by optimizing JSON schemas and increasing max tokens."
p.level = 1
p = tf.add_paragraph()
p.text = "Key Takeaways:"
p = tf.add_paragraph()
p.text = "- Never trust a single LLM for critical healthcare data. A Dual-LLM consensus engine is mandatory for safe, hallucination-free extraction."
p.level = 1


# ==========================================
# Slide 5: DEMO SUMMARY AND NEXT STEPS
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "5. DEMO SUMMARY AND NEXT STEPS"
tf = slide.placeholders[1].text_frame

p = tf.add_paragraph()
p.text = "Final Solution:"
p = tf.add_paragraph()
p.text = "- Firstsource Clinical Intelligence Hub successfully digitizes 15-minute manual chart reviews into a 10-second automated workflow."
p.level = 1
p.space_after = Pt(14)

p = tf.add_paragraph()
p.text = "Live Demo: https://medilyft-ai.onrender.com"
p = tf.add_paragraph()
p.text = "GitHub: https://github.com/rohan-9506/clinical-document-intelligence-hub"
p.space_after = Pt(14)

p = tf.add_paragraph()
p.text = "Potential Enhancements:"
p = tf.add_paragraph()
p.text = "- Direct integration with EHRs (EPIC/Cerner) via FHIR APIs."
p.level = 1
p = tf.add_paragraph()
p.text = "- Upgrading to enterprise API tiers for unlimited document lengths."
p.level = 1

prs.save("Firstsource_PoC_Deck.pptx")
print("Presentation 'Firstsource_PoC_Deck.pptx' generated successfully!")
